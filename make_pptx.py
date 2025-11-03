from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# 🎓 발표자료 초기 설정
prs = Presentation()
prs.slide_width = Inches(13.33)  # 16:9 비율
prs.slide_height = Inches(7.5)

# 🌸 색상 팔레트
PINK = RGBColor(255, 92, 138)
LIGHT_PINK = RGBColor(255, 240, 246)
GRAY = RGBColor(68, 68, 68)

titles = [
    "🌞 Smart Alarm Clock – Party Edition\nGraduation Project Presentation\n유현지 | 순천향대학교 의료IT공학과 | 2025.11.04",
    "문제 인식", "연구 목적", "시스템 구조도", "주요 기능", "기술 스택",
    "데이터 구조", "기대 효과", "디자인 컨셉",
    "시연① 일정 등록 & 알람 설정", "시연② 알람 종료 & 팝업",
    "시연③ 리포트 분석", "🌸 마무리 — 데이터로 하루를 디자인하다."
]

bodies = [
    "",
    "하루의 시작, 알람은 있지만 동기부여는 없다.\n- 일정과 연동되지 않음\n- 감성 부족\n- 피드백 부재",
    "기상 데이터를 통한 자기관리 시스템 구축\n① 일정 연동 ② 감성적 인터페이스 ③ 데이터 기반 피드백",
    "일정 입력 → 알람 모듈 → 분석 → 리포트 시각화",
    "• 일정 저장\n• 알람 설정\n• 알람 종료 팝업\n• 응원·폭죽 효과\n• 리포트 분석",
    "HTML, JS, Chart.js, LocalStorage, Python Local Server",
    "smart_alarm_routines\nsmart_alarm_logs\nsmart_alarm_analyzed",
    "하루를 데이터화하여 피드백으로 전환",
    "Cospick 감성의 글라스모픽 디자인\n핑크톤 (#ff5c8a) 기반 UI/UX",
    "일정 등록 및 알람 설정 시연 이미지 자리",
    "알람 종료 팝업 및 폭죽 효과 시연 이미지 자리",
    "Chart.js 기반 곡선형 점수 그래프 (핑크톤) 예시",
    "감성 + 데이터 기반 자기관리 시스템"
]

for i, (title, body) in enumerate(zip(titles, bodies)):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    if i <= 7:
        fill.fore_color.rgb = RGBColor(255, 255, 255)
    else:
        fill.fore_color.rgb = LIGHT_PINK

    # 제목
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(11), Inches(2))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    p = title_tf.add_paragraph()
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = PINK

    # 본문
    if body:
        body_box = slide.shapes.add_textbox(Inches(1.5), Inches(3), Inches(10.5), Inches(3.5))
        body_tf = body_box.text_frame
        body_tf.word_wrap = True
        p = body_tf.add_paragraph()
        p.text = body
        p.font.size = Pt(20)
        p.font.color.rgb = GRAY

    # 페이지 번호
    num_box = slide.shapes.add_textbox(Inches(11.5), Inches(6.7), Inches(1), Inches(0.5))
    num_tf = num_box.text_frame
    num_tf.text = f"{i+1} / 13"
    num_tf.paragraphs[0].font.size = Pt(12)
    num_tf.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)

prs.save("smart_alarm_presentation.pptx")
print("✅ smart_alarm_presentation.pptx 생성 완료!")

